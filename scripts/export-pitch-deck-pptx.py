#!/usr/bin/env python3
"""Export CurXor pitch deck to PowerPoint with speaker notes, branding, and appendix."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

try:
    from PIL import Image, ImageDraw
except ImportError:  # pragma: no cover
    Image = None  # type: ignore[assignment,misc]
    ImageDraw = None  # type: ignore[assignment,misc]

ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "docs" / "assets"
LOGO_PATH = ASSETS / "curxor-logo-mark.png"
OUTPUT = ROOT / "docs" / "CurXor-Pitch-Deck-Speaker-Notes.pptx"

# Brand tokens (curxor.ai — black + #BF5AF2 accent, monospace body on web)
BRAND = {
    "bg_dark": (12, 12, 16),
    "bg_slide": (18, 18, 22),
    "accent": (191, 90, 242),  # #BF5AF2
    "accent_soft": (160, 70, 210),
    "text_primary": (255, 255, 255),
    "text_body": (210, 210, 220),
    "text_muted": (140, 140, 155),
    "font_sans": "Segoe UI",
    "font_mono": "Consolas",
    "site": "curxor.ai",
    "version": "v1.1.0 · July 2026",
}

APPENDIX_ROWS = [
    ("Software v1.0.3", "G1/G2/G3 green · 239 smoke · 40 user flows"),
    ("Flagship Claws + Forge + Cafe", "G3 demo depth · Capital paper on camera"),
    ("Local LLM (ROCm on box)", "qwen3:8b · 38 tok/s · validated Jul 2026"),
    ("Digital bridges (Alpaca / X)", "Coded · founder-keyed until Ops Wave 1"),
    ("MS-S1 MAX validation", "Unboxed · G1–G3 closed · Act I proof plane"),
    ("Custom CurXor hardware", "Act II · post-raise"),
]

APPENDIX_DISCLAIMERS = [
    "No guaranteed returns — Capital Claw performance is strategy-dependent.",
    "Creator posts require explicit operator skill approval before bridge egress.",
    "Paper trading on camera; live trades via Alpaca on eno2 only when configured.",
]

# (slide_title, bullets, speaker_notes) — 10 slides · Fundraise D2 · Jul 2026
SLIDES: list[tuple[str, list[str], str]] = [
    (
        "CurXor",
        [
            "A sovereign AI system on your desk — with digital employees",
            "that run wealth, creation, and work on metal you own.",
            "",
            "$3M seed SAFE · 18 mo · 5 people max",
            "$3,999 once · $0/mo operate-plane API · proof-ready · pre-revenue",
        ],
        """OPENING (30–45 sec)
Hybrid hero — not personal cloud, not another chat tab.

REGISTERS
• Sovereignty investors: owned metal · no API rent · egress boundary · MS-S1 proof.
• Product investors: Capital · Creator · Work · Forge · Cafe · desk orchestrator.

METRICS
• Raise: $3M seed SAFE · 18 months · ≤5 people.
• Price: $3,999 once · $0/mo operate plane.
• Stage: G1/G2/G3 green · pre-revenue · solo founder.

CLOSE TITLE
Best next step is live LAN demo on :3080 — not a cloud sandbox.
Proof: curxor.ai/demo/investor/g3-inception-reel-v1.mp4""",
    ),
    (
        "Founder story — 20 days",
        [
            "Jun 18 — named CurXor in public · origin thread",
            "Jun 19 — first commit · Claw OS on laptop",
            "Week 1 — four flagship Claws + Forge + agent runtime",
            "Week 2 — validation hardware unboxed · Cafe · Build Plane",
            "Week 3 — v1.0.0→v1.0.3 · G1/G2/G3 · film pack",
            "",
            "182 commits · 239 smoke · 38 tok/s on validated metal",
        ],
        """J-CAL / VELOCITY (45 sec)
Specificity beats adjectives.

ARC
• Jun 18 public commitment → Jun 19 build start.
• MS-S1 is Act I validation metal — product is CurXor box + CurXor OS.
• Do not claim SpaceX/Cursor partnership.

OFFER
Inception reel (~90s) for async; live box for truth.""",
    ),
    (
        "Why now",
        [
            "Aspiration — culture says use AI; operators stall after drafts/polish",
            "Setup — always-on still means repos, keys, scripts, five tabs",
            "Product shape — chat rent · cloud agents · home cloud · light boxes",
            "Nobody ships vertical digital employees OOTB on metal you own",
            "Desk-ready compute — local inference on operator-grade UMA is real",
        ],
        """PROBLEM (45 sec)
Grammar gap is a symptom — lead holistic: want more AI, told to use it, cannot set up or stay on.

MARKET FAILURES
Cloud rent and DIY are both product-shape failures.
Wedge: appliance that closes chat → always-on team without assembly.""",
    ),
    (
        "What we build",
        [
            "CurXor box + CurXor OS (validated on MS-S1-class metal)",
            "Four pillars — compute · engine · mesh · Flight Command",
            "Ten Claws — Capital · Creator · Work · Forge · Cafe + honest previews",
            "Forge — mint custom Claws on-box · Cafe — inter-Claw proof",
            "Local cognition · egress you wire and can unplug",
        ],
        """PRODUCT (60 sec)
Say "CurXor box" first — MS-S1 is validation, not brand.

DEPTH
Flagship: Capital / Creator / Work + Forge.
Five other Claws: honest Coming Soon.
Forge + Cafe are moat — not the cold sell alone.""",
    ),
    (
        "Stack proof",
        [
            "G1 golden path · Jul 1 — on device",
            "G2 golden release · v1.0.0–v1.0.3",
            "G3 film pack · Jul 8 — inception · investor proof · desk strips",
            "239 smoke · qwen3:8b · 38 tok/s · 4.56/64 GB",
            "curxor.ai live — G3 stills · /signal · press kit",
        ],
        """PROOF (60 sec)
Lead with metrics cite card. Four pillars are systemd on box.

HONESTY
• Paper Capital on camera — say out loud.
• Live broker/social fills founder-keyed / G4 — do not claim fleet fills.
• Offer LAN demo: Flight Command :3080.

VISUALS
07-system-health-toks.png · 01-home.png · inception reel.""",
    ),
    (
        "Competition",
        [
            "Personal cloud / home server → apps, not vertical outcomes",
            "DIY agent kits → 200 hours assembly · no egress product",
            "Light assistant appliances → chat/messaging · limited compute",
            "SaaS / cloud agents → API rent · their retention",
            "",
            "CurXor = operator desk with digital employees + kill switch",
        ],
        """COMPETITIVE (45 sec)
No named competitors on the leave-behind.

CULTURAL THREAT
Free agent runtime + Mac Mini — counter is complete OS + desks + kill switch + 200 hours back.

NEVER
Abacus · partnership cosplay · breadth theater · "100x" without box data.""",
    ),
    (
        "Business model — current · evolving",
        [
            "Today: $3,999 once · $0/mo operate API · ten Claws + Forge · BYOK optional",
            "Pre-order live · pre-revenue by design (proof-first)",
            "",
            "Grows: Standard→Pro 128→Studio · Claw unlocks · Forge mint · additive bridges",
            "Hardware anchor · OTA on metal you own — not another rent trap",
        ],
        """ECONOMICS (60 sec)
Hardware is the anchor. Software evolves OTA.

REJECT
Mandatory monthly tiers as the default story (assistant-box pattern).

ACTS
Act I validates on appliance metal.
Act II CurXor-designed hardware post-investment.""",
    ),
    (
        "PMF honesty",
        [
            "Green — G1/G2/G3 · v1.0.3 · 239 smoke · G3 film pack",
            "Yellow — Flagship dogfooded · paper/sim on camera · hero Act I live",
            "Red — ARR · live fills · custom CurXor hardware (Act II)",
            "",
            "Proof-ready now · traction-ready after G4 external UAT",
        ],
        """HONESTY (45 sec)
Pre-revenue and solo founder are facts — say them.
Agent-assisted build explains velocity without inventing a team slide.
Design partners during the raise — not fake logos.""",
    ),
    (
        "Use of funds — $3M seed",
        [
            "People (eng · design · ops) — 45%",
            "Act II hardware exploration — 30%",
            "Onboarding / GTM — 10%",
            "Ops / buffer — 15%",
            "",
            "18 months · ≤5 people · prove Act I · own metal in Act II",
        ],
        """FUNDS (45 sec)
Headcount cap is intentional.
30% hardware exploration = dream-state custom box after proof, not before.
GTM buys G4 smiles and design partners — not paid vanity growth.""",
    ),
    (
        "Ask + links",
        [
            "Warm intros · live box demo on LAN :3080",
            "curxor.ai · /signal · press · inception reel · investor proof",
            "Essay + origin thread · ankur@curxor.ai · @ankurmisra",
            "",
            "No cold spam · no public raise posts · sovereign metal · one invoice",
        ],
        """CLOSE (60 sec)
Warm conversations only.
Leave inception reel for async; insist on live box when they engage.
Closing: sovereign metal · digital employees · one invoice.

LINKS
https://curxor.ai
https://curxor.ai/signal#category-film
https://curxor.ai/demo/hero-category-badge-v1.mp4
https://curxor.ai/demo/investor/g3-inception-reel-v1.mp4
https://curxor.ai/demo/investor/g3-investor-proof-v1.mp4
https://curxor.ai/press""",
    ),
]

APPENDIX_NOTES = """APPENDIX / Q&A (present only if asked)

MATURITY
• v1.0.3 appliance green — G1/G2/G3 closed Jul 8.
• Flagship desks: Capital paper, Creator queue, Work pipeline on camera.
• Bridges: coded; Ops Wave 1 keys founder-gated.
• Custom CurXor hardware: Act II post-raise.

DISCLAIMERS
• No guaranteed returns. Capital paper/live via Alpaca — strategy-dependent.
• Posts require operator approval skill — no autonomous publish.

COMMON Q
Q: Just Ollama on a PC?
A: Flight Command + ten Claws + Forge + Cafe + eno2 kill switch + OTA — appliance product.

Q: Why $3,999 vs €549 boxes?
A: Operator desks + 64GB UMA + depth — filter for buyers who run Capital/Creator/Work, not chat bots.

TIMING
Main deck ~10–12 min. Appendix 2–3 min. Q&A to 20–25."""


def rgb(name: str) -> RGBColor:
    return RGBColor(*BRAND[name])


def ensure_logo_png() -> Path:
    ASSETS.mkdir(parents=True, exist_ok=True)
    if LOGO_PATH.exists():
        return LOGO_PATH
    if Image is None or ImageDraw is None:
        raise RuntimeError("Pillow is required to generate logo asset. Run: pip install pillow")

    size = 256
    img = Image.new("RGBA", (size, size), (0, 0, 0, 255))
    draw = ImageDraw.Draw(img)
    s = size / 32
    draw.rectangle([14 * s, 6 * s, 18 * s, 10 * s], fill=(191, 90, 242, 255))
    draw.rectangle([8 * s, 14 * s, 24 * s, 16 * s], fill=(255, 255, 255, 230))
    draw.rectangle([10 * s, 18 * s, 22 * s, 20 * s], fill=(255, 255, 255, 128))
    draw.rectangle([12 * s, 22 * s, 20 * s, 24 * s], fill=(191, 90, 242, 204))
    img.save(LOGO_PATH)
    return LOGO_PATH


def set_slide_background(slide, key: str = "bg_slide") -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb(key)


def add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    font_size=18,
    bold=False,
    color_key="text_body",
    align=PP_ALIGN.LEFT,
    font_key="font_sans",
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    if not p.runs:
        p.add_run()
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color_key) if isinstance(color_key, str) else RGBColor(*color_key)
    run.font.name = BRAND[font_key]
    return tf


def add_top_accent_bar(slide) -> None:
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(0.06))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb("accent")
    bar.line.fill.background()


def add_brand_header(slide, *, logo_height=0.42) -> None:
    logo = ensure_logo_png()
    slide.shapes.add_picture(str(logo), Inches(0.55), Inches(0.18), height=Inches(logo_height))
    add_textbox(
        slide,
        Inches(1.05),
        Inches(0.16),
        Inches(2.5),
        Inches(0.35),
        "CurXor",
        font_size=16,
        bold=True,
        color_key="text_primary",
        font_key="font_mono",
    )


def add_brand_footer(slide, slide_num: int, total: int) -> None:
    add_textbox(
        slide,
        Inches(0.55),
        Inches(7.05),
        Inches(4.0),
        Inches(0.25),
        BRAND["site"],
        font_size=9,
        color_key="text_muted",
        font_key="font_mono",
    )
    add_textbox(
        slide,
        Inches(11.5),
        Inches(7.05),
        Inches(1.5),
        Inches(0.25),
        f"{slide_num} / {total}",
        font_size=9,
        color_key="text_muted",
        font_key="font_mono",
        align=PP_ALIGN.RIGHT,
    )
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(6.92), Inches(12.2), Inches(0.015))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(50, 50, 58)
    line.line.fill.background()


def apply_brand_chrome(slide, slide_num: int, total: int, *, logo_height=0.42) -> None:
    add_top_accent_bar(slide)
    add_brand_header(slide, logo_height=logo_height)
    add_brand_footer(slide, slide_num, total)


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str], notes: str, slide_num: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    apply_brand_chrome(slide, slide_num, total)

    add_textbox(
        slide,
        Inches(0.55),
        Inches(0.72),
        Inches(12.0),
        Inches(0.75),
        title,
        font_size=30,
        bold=True,
        color_key="text_primary",
    )

    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.55), Inches(1.38), Inches(0.07), Inches(0.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb("accent")
    bar.line.fill.background()

    body = slide.shapes.add_textbox(Inches(0.75), Inches(1.55), Inches(11.8), Inches(5.2))
    tf = body.text_frame
    tf.word_wrap = True
    tf.clear()

    first = True
    for line in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.text = line or " "
        p.space_after = Pt(6)
        if not p.runs:
            p.add_run()
        run = p.runs[0]
        run.font.size = Pt(10 if not line.strip() else 19)
        run.font.color.rgb = rgb("text_body")
        run.font.name = BRAND["font_sans"]

    slide.notes_slide.notes_text_frame.text = notes


def add_title_slide(prs: Presentation, bullets: list[str], notes: str, slide_num: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, "bg_dark")
    apply_brand_chrome(slide, slide_num, total, logo_height=0.55)

    add_textbox(
        slide,
        Inches(0.8),
        Inches(1.55),
        Inches(11.5),
        Inches(1.0),
        "CurXor OS",
        font_size=46,
        bold=True,
        color_key="text_primary",
        font_key="font_mono",
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(2.55),
        Inches(11.0),
        Inches(0.9),
        bullets[0],
        font_size=26,
        color_key="accent",
        font_key="font_sans",
    )

    metrics_y = 3.65
    for i, line in enumerate(bullets[2:5]):
        if not line.strip():
            continue
        add_textbox(
            slide,
            Inches(0.8),
            Inches(metrics_y + i * 0.42),
            Inches(11.0),
            Inches(0.35),
            line,
            font_size=17,
            color_key="text_body",
            font_key="font_mono",
        )

    add_textbox(
        slide,
        Inches(0.8),
        Inches(5.55),
        Inches(11.0),
        Inches(0.5),
        bullets[-1],
        font_size=15,
        color_key="text_muted",
    )
    add_textbox(
        slide,
        Inches(0.8),
        Inches(6.35),
        Inches(6.0),
        Inches(0.3),
        f"{BRAND['version']} · Sovereign Agent Appliance",
        font_size=10,
        color_key="text_muted",
        font_key="font_mono",
    )

    slide.notes_slide.notes_text_frame.text = notes


def add_appendix_slide(prs: Presentation, slide_num: int, total: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    apply_brand_chrome(slide, slide_num, total)

    add_textbox(
        slide,
        Inches(0.55),
        Inches(0.72),
        Inches(12.0),
        Inches(0.6),
        "Appendix — maturity & disclaimers",
        font_size=28,
        bold=True,
        color_key="text_primary",
    )

    rows = len(APPENDIX_ROWS) + 1
    cols = 2
    table_shape = slide.shapes.add_table(rows, cols, Inches(0.55), Inches(1.45), Inches(12.0), Inches(3.2))
    table = table_shape.table
    table.columns[0].width = Inches(7.2)
    table.columns[1].width = Inches(4.8)

    headers = ("Area", "Status")
    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.bold = True
                run.font.size = Pt(11)
                run.font.name = BRAND["font_mono"]
                run.font.color.rgb = rgb("text_primary")

    for r, (area, status) in enumerate(APPENDIX_ROWS, start=1):
        for c, value in enumerate((area, status)):
            cell = table.cell(r, c)
            cell.text = value
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.size = Pt(10)
                    run.font.name = BRAND["font_sans"]
                    run.font.color.rgb = rgb("text_body")

    y = 4.85
    add_textbox(
        slide,
        Inches(0.55),
        Inches(y),
        Inches(12.0),
        Inches(0.3),
        "Disclaimers",
        font_size=14,
        bold=True,
        color_key="accent",
    )
    for i, line in enumerate(APPENDIX_DISCLAIMERS):
        add_textbox(
            slide,
            Inches(0.55),
            Inches(y + 0.35 + i * 0.38),
            Inches(12.0),
            Inches(0.35),
            f"• {line}",
            font_size=11,
            color_key="text_muted",
        )

    slide.notes_slide.notes_text_frame.text = APPENDIX_NOTES


def build() -> Path:
    ensure_logo_png()
    total = len(SLIDES) + 1  # + appendix

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for i, (title, bullets, notes) in enumerate(SLIDES):
        slide_num = i + 1
        if i == 0:
            add_title_slide(prs, bullets, notes, slide_num, total)
        else:
            add_bullet_slide(prs, title, bullets, notes, slide_num, total)

    add_appendix_slide(prs, total, total)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT))
    return OUTPUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
